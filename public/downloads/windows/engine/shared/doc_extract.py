from __future__ import annotations

import os
import re
from dataclasses import dataclass
from typing import Any, Dict, List

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook


_WS_RE = re.compile(r"\s+")


class HwpConversionRequired(RuntimeError):
    """Raised when a .hwp document cannot be converted in-process."""


@dataclass
class DocChunk:
    chunk_id: str
    text: str
    meta: Dict[str, Any]


def _clean(text: str) -> str:
    text = text or ""
    text = text.replace("\x00", " ")
    text = _WS_RE.sub(" ", text)
    return text.strip()


def _split_text(text: str, max_chars: int) -> List[str]:
    text = text or ""
    if len(text) <= max_chars:
        return [text]

    parts: List[str] = []
    buf: List[str] = []
    size = 0
    for line in text.splitlines():
        line = (line or "").strip()
        if not line:
            continue
        if size + len(line) + 1 > max_chars and buf:
            parts.append("\n".join(buf))
            buf = [line]
            size = len(line)
        else:
            buf.append(line)
            size += len(line) + 1
    if buf:
        parts.append("\n".join(buf))

    out: List[str] = []
    for p in parts:
        if len(p) <= max_chars:
            out.append(p)
        else:
            for i in range(0, len(p), max_chars):
                out.append(p[i:i + max_chars])
    return out


def extract_chunks(path: str, *, max_chars: int = 12000, xlsx_cell_limit: int = 20000) -> List[DocChunk]:
    """Extract text chunks from local files for RAG ingestion."""
    ext = os.path.splitext(path)[1].lower().lstrip(".")

    if ext in ("txt", "md", "markdown", "log"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = (f.read() or "").strip()
        chunks: List[DocChunk] = []
        for i, part in enumerate(_split_text(raw, max_chars), start=1):
            t = _clean(part)
            if not t:
                continue
            chunks.append(DocChunk(chunk_id=f"chunk{i}", text=t, meta={"type": ext, "part": i}))
        return chunks

    if ext == "pdf":
        reader = PdfReader(path)
        chunks: List[DocChunk] = []
        for pnum, page in enumerate(reader.pages, start=1):
            try:
                txt = (page.extract_text() or "").strip()
            except Exception:
                txt = ""
            if not txt:
                continue
            for j, part in enumerate(_split_text(txt, max_chars), start=1):
                t = _clean(part)
                if not t:
                    continue
                chunks.append(DocChunk(chunk_id=f"p{pnum}-c{j}", text=t, meta={"type": "pdf", "page": pnum, "chunk": j}))
        return chunks

    if ext == "docx":
        doc = DocxDocument(path)
        paras = [(p.text or "").strip() for p in doc.paragraphs]
        paras = [p for p in paras if p]
        buf = "\n".join(paras).strip()
        chunks: List[DocChunk] = []
        for i, part in enumerate(_split_text(buf, max_chars), start=1):
            t = _clean(part)
            if not t:
                continue
            chunks.append(DocChunk(chunk_id=f"chunk{i}", text=t, meta={"type": "docx", "chunk": i}))
        return chunks

    if ext == "pptx":
        pres = Presentation(path)
        chunks: List[DocChunk] = []
        for sidx, slide in enumerate(pres.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            body = "\n".join(texts).strip()
            if not body:
                continue
            for j, part in enumerate(_split_text(body, max_chars), start=1):
                t = _clean(part)
                if not t:
                    continue
                chunks.append(DocChunk(chunk_id=f"s{sidx}-c{j}", text=t, meta={"type": "pptx", "slide": sidx, "chunk": j}))
        return chunks

    if ext in ("xlsx", "xlsm", "xltx", "xltm"):
        wb = load_workbook(path, read_only=True, data_only=True)
        chunks: List[DocChunk] = []
        for name in wb.sheetnames:
            ws = wb[name]
            lines: List[str] = []
            cells = 0
            for row in ws.iter_rows(values_only=True):
                row_vals = ["" if v is None else str(v) for v in row]
                line = "\t".join(row_vals).strip()
                if line:
                    lines.append(line)
                cells += len(row_vals)
                if cells >= xlsx_cell_limit:
                    break
            body = "\n".join(lines).strip()
            if not body:
                continue
            for j, part in enumerate(_split_text(body, max_chars), start=1):
                t = _clean(part)
                if not t:
                    continue
                chunks.append(DocChunk(chunk_id=f"{name}-c{j}", text=t, meta={"type": "xlsx", "sheet": name, "chunk": j}))
        return chunks

    if ext == "hwp":
        raise HwpConversionRequired("HWP requires external conversion to text or PDF before ingest.")

    raise ValueError(f"Unsupported file type: {ext}")
